"""learning-kb 测试夹具与材料工厂。

所有测试要求:在 learning-kb/.venv 下运行
(依赖已由 scripts/setup.ps1 安装,并额外安装 pytest)。
"""
import json
import math
import struct
import sys
import wave
from pathlib import Path

import pytest

SKILL = Path(__file__).resolve().parent.parent
SCRIPTS = SKILL / "scripts"
sys.path.insert(0, str(SCRIPTS))

import common  # noqa: E402


# ---------------------------------------------------------------- 基础夹具

@pytest.fixture
def state(tmp_path):
    """空的 .learning-kb 状态目录(带 vault 父目录),返回 str 路径。"""
    vault = tmp_path / "vault"
    vault.mkdir()
    state_dir = vault / ".learning-kb"
    state_dir.mkdir()
    return str(state_dir)


@pytest.fixture
def state_with_brief(state):
    common.write_json(Path(state) / "brief.json", {
        "version": 1,
        "vault_path": str(Path(state).parent),
        "state_dir": state,
        "goal": "测试目标",
        "external_policy": "official",
        "vision": "no",
    })
    return state


@pytest.fixture
def materials(tmp_path):
    """材料工厂目录,返回路径。"""
    m = tmp_path / "materials"
    m.mkdir()
    return m


def load_sources(state_dir):
    p = Path(state_dir) / "sources.json"
    if not p.exists():
        return []
    return common.read_json(p)["sources"]


def get_source(state_dir, sid):
    return next(s for s in load_sources(state_dir) if s["id"] == sid)


# ---------------------------------------------------------------- 材料工厂

def make_png(path: Path, width=40, height=30):
    import pymupdf as fitz
    pix = fitz.Pixmap(fitz.csRGB, fitz.IRect(0, 0, width, height))
    pix.clear_with(255)
    for y in range(height // 2):
        for x in range(width):
            pix.set_pixel(x, y, (200, 30, 30))
    pix.save(str(path))
    return path


def make_pdf(path: Path, pages, image_every=0):
    """pages: [(title, body), ...];image_every>0 时每 N 页插一张图。"""
    import pymupdf as fitz
    doc = fitz.open()
    for i, (title, body) in enumerate(pages, 1):
        page = doc.new_page(width=595, height=842)
        page.insert_text((72, 100), title, fontsize=18, fontname="china-s")
        page.insert_text((72, 160), body, fontsize=11, fontname="china-s")
        if image_every and i % image_every == 0:
            img = path.parent / f"_tmp_{i}.png"
            make_png(img)
            page.insert_image(fitz.Rect(72, 300, 172, 380), filename=str(img))
            img.unlink()
    doc.save(str(path))
    return path


def make_docx(path: Path, heading=None, body=None, table=None, image=False):
    import docx
    d = docx.Document()
    if heading:
        d.add_heading(heading, level=1)
    if body:
        d.add_paragraph(body)
    if table:
        t = d.add_table(rows=len(table), cols=len(table[0]))
        for i, row in enumerate(table):
            for j, cell in enumerate(row):
                t.cell(i, j).text = cell
    if image:
        img = path.parent / "_tmp_docx.png"
        make_png(img)
        d.add_picture(str(img))
        img.unlink()
    d.save(str(path))
    return path


def make_pptx(path: Path, slides, notes=None):
    """slides: [(title, body), ...]; notes: 每页备注,可 None。"""
    from pptx import Presentation
    prs = Presentation()
    for i, (title, body) in enumerate(slides):
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = title
        slide.placeholders[1].text = body
        if notes:
            slide.notes_slide.notes_text_frame.text = notes[i]
    prs.save(str(path))
    return path


def make_epub(path: Path, chapters, image=False):
    import ebooklib
    from ebooklib import epub
    book = epub.EpubBook()
    book.set_identifier("test-001")
    book.set_title("Test Book")
    book.set_language("zh")
    items = []
    for i, (title, body) in enumerate(chapters, 1):
        html = f"<h1>{title}</h1><p>{body}</p>"
        if image and i == 1:
            img = path.parent / "_tmp_epub.png"
            make_png(img)
            data = img.read_bytes()
            img.unlink()
            item = epub.EpubItem(uid=f"img{i}", file_name=f"img{i}.png",
                                 media_type="image/png", content=data)
            book.add_item(item)
            html += f'<p><img src="img{i}.png" alt="test image"/></p>'
        c = epub.EpubHtml(title=title, file_name=f"chap{i}.xhtml", lang="zh")
        c.content = html
        book.add_item(c)
        items.append(c)
    book.toc = tuple(items)
    book.add_item(epub.EpubNcx())
    book.add_item(epub.EpubNav())
    book.spine = ["nav"] + items
    epub.write_epub(str(path), book)
    return path


def make_wav(path: Path, seconds=1.0, freq=440.0):
    rate = 16000
    n = int(seconds * rate)
    with wave.open(str(path), "wb") as w:
        w.setnchannels(1)
        w.setsampwidth(2)
        w.setframerate(rate)
        frames = b"".join(
            struct.pack("<h", int(3000 * math.sin(2 * math.pi * freq * t / rate)))
            for t in range(n))
        w.writeframes(frames)
    return path


def make_chat(path: Path, lines):
    path.write_text("\n\n".join(lines), encoding="utf-8")
    return path


# ---------------------------------------------------------------- 模拟 LLM 阶段

def write_proposal(state_dir, src_id, items):
    p = Path(state_dir) / "ir" / "proposals" / f"{src_id}.json"
    common.write_json(p, {"version": 1, "src_id": src_id, "items": items})
    common.update_source(state_dir, src_id, status="proposed")
    return p


def merge_proposals(state_dir):
    """把全部提案合并进 knowledge.json(测试用简化 merge:
    同 id 并入 sources,新 id 追加)。返回 knowledge items。"""
    base = Path(state_dir)
    kf = base / "ir" / "knowledge.json"
    knowledge = common.read_json(kf) if kf.exists() else {"version": 1, "items": []}
    by_id = {i["id"]: i for i in knowledge["items"]}
    for prop_file in sorted((base / "ir" / "proposals").glob("*.json")):
        prop = common.read_json(prop_file)
        for item in prop["items"]:
            if item["id"] in by_id:
                existing = by_id[item["id"]]
                changed = False
                for src in item.get("sources", []):
                    if src not in existing["sources"]:
                        existing["sources"].append(src)
                        changed = True
                if item.get("detail") and item["detail"] != existing.get("detail"):
                    existing["detail"] = item["detail"]
                    changed = True
                if changed:
                    existing["status"] = "updated"
            else:
                item["status"] = "accepted"
                knowledge["items"].append(item)
                by_id[item["id"]] = item
    common.write_json(kf, knowledge)
    return knowledge["items"]


def compile_notes(state_dir, topics, bodies):
    """topics: [{topic, path, items}]; bodies: {item_id: markdown 正文}。
    写 staging 笔记 + Start Here(含 managed 区域与互链)。"""
    base = Path(state_dir)
    staging = base / "staging"
    common.write_json(base / "ir" / "topics.json", topics)
    notes = []
    for t in topics:
        sections = []
        for iid in t["items"]:
            sections.append(f"## {bodies[iid]['title']}\n\n{bodies[iid]['detail']}\n")
        links = "".join(
            f"- [[{other['path']}]]\n"
            for other in topics if other["path"] != t["path"])
        content = (f"<!-- lkb-managed:start -->\n# {t['topic']}\n\n"
                   + "\n".join(sections)
                   + f"\n## 相关\n\n{links}<!-- lkb-managed:end -->\n")
        note = staging / (t["path"] + ".md")
        note.parent.mkdir(parents=True, exist_ok=True)
        note.write_text(
            f'---\ntitle: "{t["topic"]}"\ntags: [t]\nlkb-generated: true\n---\n\n' + content,
            encoding="utf-8")
        notes.append(note)
    top = next((t for t in topics if "/" not in t["path"]), None)
    items = "".join(
        f"- [[{t['path']}]]\n"
        for t in topics if (top is None or t["path"] != top["path"]))
    head = f"- [[{top['path']}]]\n" if top else ""
    (staging / "Start Here.md").write_text(
        f'---\ntitle: 开始阅读\ntags: [moc]\nlkb-generated: true\n---\n\n'
        f'<!-- lkb-managed:start -->\n# 开始阅读\n\n## 主题地图\n\n{head}{items}<!-- lkb-managed:end -->\n',
        encoding="utf-8")
    return staging


def register_md(state_dir, md_path):
    """markdown/txt 手工登记(模拟 kb-ingest 协议)。"""
    entry, dup = common.register(state_dir, "md", md_path)
    if dup:
        return entry, True
    text = Path(md_path).read_text(encoding="utf-8")
    out = f"<!-- src:{entry['id']} title:{entry['title']} -->\n\n{text}"
    (Path(state_dir) / "extracted" / f"{entry['id']}.md").write_text(out, encoding="utf-8")
    common.update_source(state_dir, entry["id"],
                         extracted_file=f"extracted/{entry['id']}.md", status="parsed")
    return entry, False
